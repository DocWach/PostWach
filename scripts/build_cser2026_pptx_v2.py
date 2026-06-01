"""Build CSER 2026 presentation v2 matching the restructured paper outline."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Paths ---
BASE = r"C:\Users\pfwac\OneDrive - University of Arizona\Documents"
FIG_DIR = os.path.join(
    BASE, r"03 Projects\00 My Research\01 PostWach\Papers\SE_Math_Foundations"
    r"\isomorphism-library\paper\figures"
)
OUT_PATH = os.path.join(
    BASE, r"02 My Outreach\CSER 2026 - Morphisms"
    r"\Wach_Sandman_Iyer_CSER2026_Presentation_v2.pptx"
)

# --- Design constants ---
NAVY = RGBColor(0x00, 0x33, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
FONT_NAME = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_BAR_H = Inches(1.1)
BODY_TOP = Inches(1.3)
BODY_LEFT = Inches(0.6)
BODY_WIDTH = Inches(12.1)
BODY_HEIGHT = Inches(5.8)


def fig_path(name):
    p = os.path.join(FIG_DIR, name)
    return p if os.path.exists(p) else None


# --- Helpers ---

def add_title_bar(slide, title_text):
    """Add a navy title bar at the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, TITLE_BAR_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].add_run()
    run.text = title_text
    run.font.size = Pt(32)
    run.font.color.rgb = WHITE
    run.font.name = FONT_NAME
    run.font.bold = True

    # Vertical centering via margins
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)


def add_slide_number(slide, num):
    """Add slide number in bottom-right corner."""
    txBox = slide.shapes.add_textbox(
        Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(11)
    run.font.color.rgb = NAVY
    run.font.name = FONT_NAME


def add_body_text(slide, lines, top=None, left=None, width=None, height=None,
                  font_size=18, bold_first=False, line_spacing=1.15):
    """Add a text box with bullet lines."""
    t = top or BODY_TOP
    l = left or BODY_LEFT
    w = width or BODY_WIDTH
    h = height or BODY_HEIGHT

    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)

        # Handle indented sub-bullets
        indent_level = 0
        text = line
        while text.startswith("  "):
            indent_level += 1
            text = text[2:]
        text = text.lstrip("- ").lstrip("\u2022 ")

        if indent_level > 0:
            p.level = indent_level

        run = p.add_run()
        run.text = "\u2022  " + text if indent_level == 0 else "\u2013  " + text
        run.font.size = Pt(font_size - (indent_level * 2))
        run.font.color.rgb = BLACK
        run.font.name = FONT_NAME
        if bold_first and i == 0:
            run.font.bold = True


def add_table(slide, headers, rows, left=None, top=None, width=None,
              col_widths=None, font_size=14):
    """Add a formatted table with navy header row and alternating shading."""
    l = left or BODY_LEFT
    t = top or Inches(1.5)
    w = width or BODY_WIDTH
    n_rows = len(rows) + 1
    n_cols = len(headers)
    row_h = Inches(0.4)
    tbl_h = row_h * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, l, t, w, tbl_h)
    table = table_shape.table

    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.color.rgb = WHITE
                run.font.name = FONT_NAME
                run.font.bold = True

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            # Alternating shading
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = BLACK
                    run.font.name = FONT_NAME

    return table_shape


def add_image_or_placeholder(slide, img_name, left, top, width, height=None):
    """Add an image if it exists; otherwise a placeholder text box."""
    p = fig_path(img_name)
    if p:
        if height:
            slide.shapes.add_picture(p, left, top, width, height)
        else:
            slide.shapes.add_picture(p, left, top, width)
    else:
        txBox = slide.shapes.add_textbox(left, top, width, height or Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = f"[Figure: {img_name}]"
        run.font.size = Pt(14)
        run.font.color.rgb = MID_GRAY
        run.font.name = FONT_NAME
        run.font.italic = True
        # Add border
        shape = txBox
        shape.line.color.rgb = MID_GRAY
        shape.line.width = Pt(1)


# ============================================================
# BUILD PRESENTATION
# ============================================================

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank

slide_num = 0

# ------------------------------------------------------------------
# SLIDE 1: Title
# ------------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)

# Full navy background
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = NAVY
bg_shape.line.fill.background()

# Title text
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Toward a Library of Isomorphic Patterns\nfor Systems Engineering"
run.font.size = Pt(40)
run.font.color.rgb = WHITE
run.font.name = FONT_NAME
run.font.bold = True

# Authors
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(1))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Paul F. Wach (University of Arizona)"
run2.font.size = Pt(24)
run2.font.color.rgb = WHITE
run2.font.name = FONT_NAME

p3 = tf2.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "Brandt Sandmann, Adithya Iyer (Virginia Tech)"
run3.font.size = Pt(24)
run3.font.color.rgb = WHITE
run3.font.name = FONT_NAME

# Conference
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.3), Inches(0.8))
tf3 = txBox3.text_frame
p4 = tf3.paragraphs[0]
p4.alignment = PP_ALIGN.CENTER
run4 = p4.add_run()
run4.text = "CSER 2026  |  Arlington, VA"
run4.font.size = Pt(20)
run4.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
run4.font.name = FONT_NAME

# ------------------------------------------------------------------
# SLIDE 2: Motivation
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Motivation")
add_slide_number(slide, slide_num + 1)
add_body_text(slide, [
    "Cross-domain analogies (mechanical, electrical, hydraulic, acoustic) are "
    "prevalent in engineering but have never been formally characterized "
    "using SE-theoretic formalism",
    "6 established pairs across 5+ physical domains identified in the literature",
    "No formal characterization of how morphism type changes with abstraction "
    "level or numerical method",
    "If SE is to be domain-agnostic, its formalisms must characterize "
    "equivalences across domains, not just within them",
    "Gap: the SE community uses morphism language informally; "
    "no library of formally characterized patterns exists",
], font_size=20)

# ------------------------------------------------------------------
# SLIDE 3: Contribution
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Contribution")
add_slide_number(slide, slide_num + 1)
add_body_text(slide, [
    "A repeatable six-step library methodology for formally characterizing "
    "cross-domain morphic patterns using Wymore's T3SD formalism",
    "Two library entries demonstrated:",
    "  Simple: electrical capacitor and hydraulic micro-bubble (first-order, "
    "single-element)",
    "  Complex: mass-spring-damper and series RLC circuit (second-order, "
    "multi-element)",
    "Isomorphic degradation documented on two independent axes:",
    "  Structural (D_s): from abstraction and modeling choices",
    "  Behavioral (D_b): from discretization and numerical method",
    "Key insight: isomorphisms are nominal reference points in the physics; "
    "degree of homomorphism is lifecycle-continuous",
], font_size=19)

# ------------------------------------------------------------------
# SLIDE 4: Background, T3SD Formalism
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Background: T3SD Formalism")
add_slide_number(slide, slide_num + 1)
add_body_text(slide, [
    "System model as Wymore five-tuple:  Z = (S\u2082, I\u2082, O\u2082, "
    "N\u2082, R\u2082)",
    "  S = state space,  I = input space,  O = output space",
    "  N = next-state function,  R = readout function",
    "Homomorphism conditions (i-v):",
    "  (i) Surjection on inputs,  (ii) Surjection on outputs,  "
    "(iii) Surjection on states",
    "  (iv) Commutativity of N (next-state),  (v) Commutativity of R (readout)",
    "Isomorphism = bijective homomorphism (all surjections become bijections)",
    "T3SD separates specification from execution; this distinction produces "
    "the D_s / D_b split",
], font_size=19)

# ------------------------------------------------------------------
# SLIDE 5: Degree of Homomorphism
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Degree of Homomorphism  D\u2095")
add_slide_number(slide, slide_num + 1)
add_body_text(slide, [
    "D\u2095 = (D_s, D_b)",
    "D_s = structural degree (analytical, from the Z tuple):",
    "  Presented as individual ratios: D_s = (DoH_S, DoH_I, DoH_O)",
    "  Each ratio: reciprocal of preimage count, averaged across elements",
    "  Scalar average provided as convenience; individual ratios carry more "
    "information",
    "D_b = behavioral degree (empirical, from verification activity):",
    "  D_b = max\u209C |y\u2081(t) \u2212 y\u2082(t)| under the variable mapping",
    "D_s and D_b are complementary and independent:",
    "  High D_s with D_b > 0: structural fidelity degraded by implementation",
    "  Low D_s with D_b = 0: structural info loss without behavioral consequence",
    "For exact morphisms: N, R commutativity guaranteed by conditions iv-v",
], font_size=18)

# ------------------------------------------------------------------
# SLIDE 6: Well-Formed System Models (Step 0)
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Methodology: Well-Formed System Models (Step 0)")
add_slide_number(slide, slide_num + 1)
add_body_text(slide, [
    "Prerequisite: distinguish constitutive properties from inputs and states",
    "Constitutive properties (k, c, m, L, R, C) define N\u2082 (system identity):",
    "  Changing a constitutive property = different system, not same system "
    "with different input",
    "Inputs (F(t), E(t), p(t)) = external excitation applied to the system",
    "States (x, dx/dt, q, dq/dt) = internal variables evolving under N\u2082",
    "Analogy convention is also a modeling choice:",
    "  Force-voltage (impedance) vs. force-current (mobility)",
    "  Each pairs with exactly one circuit topology (series vs. parallel)",
    "  This paper uses force-voltage throughout",
    "Reference: Wach and Salado (2022), modeling choices affect "
    "problem space definition",
], font_size=18)

# ------------------------------------------------------------------
# SLIDE 7: Six-Step Procedure
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Six-Step Characterization Procedure")
add_slide_number(slide, slide_num + 1)

# Steps as text
add_body_text(slide, [
    "Step 0: Establish well-formed system models (constitutive properties, "
    "convention)",
    "Step 1: Define the system pair",
    "Step 2: Model each system at multiple abstraction levels",
    "Step 3: Characterize the morphism at each level (conditions i-v, D\u2095)",
    "Step 4: Simulate and verify output equivalence (D_b = 0?)",
    "Step 5: Discretize and measure degradation (multiple numerical methods)",
    "Step 6: Compute D\u2095 = (D_s, D_b)",
], top=BODY_TOP, height=Inches(3.0), font_size=20)

# Catalog table below
catalog_headers = ["#", "System Pair", "Domains", "Source"]
catalog_rows = [
    ["1", "MSD and RLC Circuit", "Mechanical / Electrical", "Olson, Firestone, Karnopp"],
    ["2", "Translational and Rotational Mech.", "Mechanical / Mechanical", "Karnopp"],
    ["3", "Electrical and Hydraulic", "Electrical / Hydraulic", "Olson, Karnopp"],
    ["4", "Electrical and Acoustic", "Electrical / Acoustic", "Olson"],
    ["5", "Electrical and Magnetic", "Electrical / Magnetic", "Karnopp"],
    ["6", "Mechanical and Fluid (compressible)", "Mechanical / Fluid", "Karnopp"],
]
add_table(slide, catalog_headers, catalog_rows,
          top=Inches(4.6), font_size=13,
          col_widths=[Inches(0.5), Inches(4.0), Inches(3.8), Inches(3.8)])

# ------------------------------------------------------------------
# SLIDE 8: Entry 1, Capacitor and Bubble
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Entry 1: Electrical Capacitor and Hydraulic Micro-Bubble")
add_slide_number(slide, slide_num + 1)

# Left column: text
add_body_text(slide, [
    "First-order ODEs:",
    "  Capacitor: I(t) = C dV/dt",
    "  Bubble: \u0394Q(t) = C\u2095 dp/dt",
    "Variable mapping:",
    "  V \u2194 p  (voltage, pressure)",
    "  I \u2194 \u0394Q  (current, flow rate)",
    "  q \u2194 \u0394v  (charge, volume change)",
    "  C \u2194 \u03B2  (capacitance, compliance)",
    "2-state isomorphism: D_s = (1.0, 1.0, 1.0)",
    "D\u2095 = ((1.0, 1.0, 1.0), 0)",
    "Simplest possible case: clean bijective mapping, no compositional "
    "complexity",
], left=BODY_LEFT, width=Inches(6.5), font_size=17)

# Right: figure
add_image_or_placeholder(
    slide, "capacitor_bubble_step_response.png",
    Inches(7.5), Inches(1.5), Inches(5.2), Inches(3.8)
)

# ------------------------------------------------------------------
# SLIDE 9: Entry 2, MSD and Series RLC
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Entry 2: Mass-Spring-Damper and Series RLC Circuit")
add_slide_number(slide, slide_num + 1)

# Left text
add_body_text(slide, [
    "Second-order ODEs, force-voltage convention:",
    "  MSD: m x'' + c x' + k x = F(t)",
    "  RLC: L q'' + R q' + (1/C) q = E(t)",
    "Variable mapping: x \u2194 q, v \u2194 i, F \u2194 E, m \u2194 L, "
    "c \u2194 R, k \u2194 1/C",
    "2-state isomorphism: D_s = (1.0, 1.0, 1.0)",
    "4-to-2-state homomorphism: D_s = (0.67, 1.0, 0.67)",
    "  (four MSD states mapped to two RLC states)",
    "D\u2095 at physics level: ((1.0, 1.0, 1.0), 0)",
    "Key: abstraction level changes the morphism type",
], left=BODY_LEFT, width=Inches(6.5), font_size=17)

# Right: step response figure
add_image_or_placeholder(
    slide, "step_response_overlay.png",
    Inches(7.5), Inches(1.5), Inches(5.2), Inches(3.8)
)

# Diagrams at bottom if space
add_image_or_placeholder(
    slide, "msd_system_diagram.png",
    Inches(7.5), Inches(5.5), Inches(2.5), Inches(1.6)
)
add_image_or_placeholder(
    slide, "rlc_circuit_diagram.png",
    Inches(10.2), Inches(5.5), Inches(2.5), Inches(1.6)
)

# ------------------------------------------------------------------
# SLIDE 10: Discretization Degradation
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Discretization Degradation")
add_slide_number(slide, slide_num + 1)

# Left: text + table
add_body_text(slide, [
    "Numerical method determines whether physics-level isomorphism "
    "survives implementation",
    "Pattern holds across system orders (first-order and second-order pairs)",
], left=BODY_LEFT, width=Inches(6.0), font_size=19, top=BODY_TOP,
              height=Inches(1.5))

disc_headers = ["Method", "D_b Scaling", "Preserves Isomorphism?"]
disc_rows = [
    ["Euler (forward)", "O(\u0394t)", "No"],
    ["RK4", "O(\u0394t\u2074)", "No (small error)"],
    ["Exact (analytical)", "0", "Yes"],
]
add_table(slide, disc_headers, disc_rows,
          left=BODY_LEFT, top=Inches(3.0), width=Inches(6.0), font_size=15)

# Right: figure
add_image_or_placeholder(
    slide, "discretization_degradation.png",
    Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.0)
)

# ------------------------------------------------------------------
# SLIDE 11: Cross-Entry Comparison
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Cross-Entry Comparison")
add_slide_number(slide, slide_num + 1)

comp_headers = ["Configuration", "D_s", "D_b", "D\u2095"]
comp_rows = [
    ["Cap/Bubble, 2-state iso", "(1.0, 1.0, 1.0)", "0",
     "((1.0, 1.0, 1.0), 0)"],
    ["MSD/RLC, 2-state iso", "(1.0, 1.0, 1.0)", "0",
     "((1.0, 1.0, 1.0), 0)"],
    ["MSD/RLC, 4\u21922 homo", "(0.67, 1.0, 0.67)", "0",
     "((0.67, 1.0, 0.67), 0)"],
    ["MSD/RLC, Euler \u0394t=0.5", "(1.0, 1.0, 1.0)", "O(10\u2070)",
     "((1.0, 1.0, 1.0), O(10\u2070))"],
]
add_table(slide, comp_headers, comp_rows,
          top=Inches(1.5), font_size=15,
          col_widths=[Inches(3.5), Inches(2.8), Inches(1.8), Inches(4.0)])

add_body_text(slide, [
    "What the complex case reveals that the simple cannot:",
    "  Abstraction-level dependence: morphism type changes with granularity",
    "  Modeling choice sensitivity: constitutive properties define system "
    "identity",
    "  Compositional questions: do element-level isomorphisms compose to "
    "system-level?",
    "  Multiple valid morphisms through different analogy conventions",
], top=Inches(4.3), font_size=17)

# ------------------------------------------------------------------
# SLIDE 12: Discussion, Verification Connection
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Discussion: Verification Connection")
add_slide_number(slide, slide_num + 1)
add_body_text(slide, [
    "Degree of homomorphism is computable at any point in the lifecycle, "
    "not a one-time certification",
    "As physical systems age (parameter drift, component degradation), "
    "the morphic relationship between model and real system degrades",
    "Digital twin fidelity is measurable through D\u2095 between "
    "twin model and physical system",
    "Verification planning: required D\u2095 between system design "
    "and verification model affects resources needed",
    "Isomorphisms exist in the physics independently of models, but "
    "the characterization depends on modeling choices:",
    "  Abstraction level, discretization method, and analogy convention "
    "all affect whether a system-level isomorphism is reflected in the "
    "model-level representation",
    "D_s and D_b independence demonstrated: structural info loss "
    "without behavioral consequence (4\u21922 state), and behavioral "
    "divergence without structural loss (Euler)",
], font_size=18)

# ------------------------------------------------------------------
# SLIDE 13: Future Work
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Future Work")
add_slide_number(slide, slide_num + 1)
add_body_text(slide, [
    "4 remaining established pairs (translational/rotational, "
    "electrical/acoustic, electrical/magnetic, mechanical/fluid)",
    "Configuration space: up to 8 distinct pairings for MSD/RLC "
    "(2 mechanical \u00d7 2 electrical \u00d7 2 conventions)",
    "Compositional morphisms: do system-level isomorphisms decompose "
    "to element-level? Transitivity?",
    "D\u2095 evolution: extending D_s to incorporate N and R for "
    "approximate morphisms",
    "DEVS reformulation: embeds simulation semantics, eliminates "
    "specification-simulation gap, native compositional support",
    "Category-theoretic framework: systems as objects, morphisms as "
    "arrows, discretization as functor",
    "Nonlinear and stochastic extensions",
    "Problem space implications of convention and modeling choices",
], font_size=18)

# ------------------------------------------------------------------
# SLIDE 14: Conclusion
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Conclusion")
add_slide_number(slide, slide_num + 1)
add_body_text(slide, [
    "Methodology: a repeatable six-step procedure for formally "
    "characterizing cross-domain morphic patterns using Wymore's T3SD",
    "Two library entries spanning three physical domains:",
    "  Simple (capacitor/bubble): clean bijective mapping, first-order",
    "  Complex (MSD/RLC): abstraction-dependent degradation, second-order",
    "Degree of homomorphism D\u2095 = (D_s, D_b) captures structural "
    "correspondence (analytical) and behavioral fidelity (empirical)",
    "Isomorphic degradation on two axes:",
    "  Structural: D_s changes from (1.0, 1.0, 1.0) to (0.67, 1.0, 0.67)",
    "  Behavioral: D_b scales as O(\u0394t) for Euler, O(\u0394t\u2074) "
    "for RK4, preserved by exact methods",
    "Isomorphisms are nominal reference points; degree of homomorphism "
    "is lifecycle-continuous",
    "A systematically built library would yield a foundational reference "
    "for domain-agnostic SE",
], font_size=18)

# ------------------------------------------------------------------
# SLIDE 15: Questions
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)

# Full navy background
bg2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
)
bg2.fill.solid()
bg2.fill.fore_color.rgb = NAVY
bg2.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Questions?"
run.font.size = Pt(48)
run.font.color.rgb = WHITE
run.font.name = FONT_NAME
run.font.bold = True

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(2.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True

p1 = tf2.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
run1 = p1.add_run()
run1.text = "Paul F. Wach"
run1.font.size = Pt(24)
run1.font.color.rgb = WHITE
run1.font.name = FONT_NAME

p2 = tf2.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "paulwach@arizona.edu"
run2.font.size = Pt(20)
run2.font.color.rgb = RGBColor(0xAA, 0xCC, 0xEE)
run2.font.name = FONT_NAME

p3 = tf2.add_paragraph()
p3.space_before = Pt(20)
p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run()
run3.text = "This work was supported by SERC project WRT-2516,\n" \
            '"Systems Engineering Beyond the Horizon"'
run3.font.size = Pt(16)
run3.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
run3.font.name = FONT_NAME
run3.font.italic = True

add_slide_number(slide, slide_num + 1)

# ------------------------------------------------------------------
# BACKUP SLIDE 1: D_h Calculation Detail
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Backup: Degree of Homomorphism Calculation Detail")
add_slide_number(slide, slide_num + 1)

add_body_text(slide, [
    "4-to-2 state mapping (MSD to RLC):",
    "  MSD states: {(x, v), (x, 0), (0, v), (0, 0)}  (4 states)",
    "  RLC states: {(q, i), (q_0, 0)}  (2 states, simplified)",
    "DoH_S calculation:",
    "  Each RLC state has 2 MSD preimages on average",
    "  DoH_S = avg(1/|preimage|) = avg(1/2) = 0.5 ... but with the "
    "actual mapping: (1/2 + 1/2 + 1/1 + ...)/4",
    "  Result: DoH_S = 0.67",
    "DoH_I = 1.0 (inputs remain bijective: F \u2194 E)",
    "DoH_O = 0.67 (outputs follow the state-space reduction pattern)",
    "D_s = (0.67, 1.0, 0.67), scalar average = 0.78",
    "Individual ratios reveal that input space is fully preserved while "
    "state and output spaces lose information",
], font_size=17)

# ------------------------------------------------------------------
# BACKUP SLIDE 2: Force-Voltage vs. Force-Current
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Backup: Force-Voltage vs. Force-Current Conventions")
add_slide_number(slide, slide_num + 1)

add_body_text(slide, [
    "Two established analogy conventions for mechanical-electrical pairs:",
    "  Force-voltage (impedance): force \u2194 voltage, velocity \u2194 "
    "current (Olson, Maxwell)",
    "  Force-current (mobility): force \u2194 current, velocity \u2194 "
    "voltage (Firestone)",
    "Each convention pairs with exactly one circuit topology:",
    "  Force-voltage \u2194 series RLC",
    "  Force-current \u2194 parallel RLC",
    "  Cross-pairings fail (morphism conditions violated)",
    "Convention and topology are coupled, reflecting underlying "
    "series-parallel duality",
    "4-pairing analysis: 2 conventions \u00d7 2 topologies = 4; "
    "only 2 of 4 produce valid isomorphisms",
    "Full 8-configuration analysis (adding Maxwell viscoelastic MSD) "
    "is future work",
], font_size=17)

# ------------------------------------------------------------------
# BACKUP SLIDE 3: Parameter Values
# ------------------------------------------------------------------
slide_num += 1
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Backup: Parameter Values Used")
add_slide_number(slide, slide_num + 1)

param_headers = ["Parameter", "MSD Value", "RLC Value", "Unit Mapping"]
param_rows = [
    ["Mass / Inductance", "m = 1.0 kg", "L = 1.0 H", "kg \u2194 H"],
    ["Damping / Resistance", "c = 0.5 N\u00b7s/m", "R = 0.5 \u03a9",
     "N\u00b7s/m \u2194 \u03a9"],
    ["Stiffness / Inv. Cap.", "k = 4.0 N/m", "1/C = 4.0 1/F",
     "N/m \u2194 1/F"],
    ["Force / Voltage", "F = 1.0 N (step)", "E = 1.0 V (step)", "N \u2194 V"],
]
add_table(slide, param_headers, param_rows,
          top=Inches(1.5), font_size=15,
          col_widths=[Inches(3.0), Inches(3.0), Inches(3.0), Inches(3.1)])

param2_headers = ["Parameter", "Capacitor Value", "Bubble Value", "Unit Mapping"]
param2_rows = [
    ["Capacitance / Compliance", "C = 1.0 F", "\u03b2 = 1.0 m\u00b3/Pa",
     "F \u2194 m\u00b3/Pa"],
    ["Current / Flow rate", "I = 1.0 A (step)", "\u0394Q = 1.0 m\u00b3/s (step)",
     "A \u2194 m\u00b3/s"],
]
add_table(slide, param2_headers, param2_rows,
          top=Inches(4.2), font_size=15,
          col_widths=[Inches(3.0), Inches(3.0), Inches(3.0), Inches(3.1)])

# ============================================================
# SAVE
# ============================================================
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
